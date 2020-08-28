from pptx import Presentation
from pptx.util import Inches, Pt

from PPTExcelUtils import TableChartMaker, pptText


def alter_table_fontsize(table, size):
    """
    修改表格字体大小
    :param table: 表格
    :param size: 字体大小
    :return:
    """
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)


def fill_table(rows_nums, cols_nums, table, df):
    """
    填充表格数据
    :param rows_nums:  行数
    :param cols_nums:  列数
    :param table:      table对象
    :param df:         数据df
    :return:
    """
    for i in range(rows_nums):
        for j in range(cols_nums):
            table.cell(i, j).text = str(df.iloc[i, j])
            # 填充表格颜色
            # table.cell(i, j).fill.solid()
            # table.cell(i, j).fill.fore_color.rgb = RGBColor(255, 99, 71)


class PPT(object):
    def __init__(self):
        self.read_path = 'D:/项目文档/分析报告/model.pptx'
        self.save_path = "D:/项目文档/分析报告/1www.pptx"
        self.pic_path = "D:/项目文档/分析报告/pic/"

    def p2_operate(self):
        prs = Presentation(self.read_path)
        slide = prs.slides[1]
        body_shapes = slide.shapes
        # 获取文本所需数据
        df_profit, df_rate = TableChartMaker.Sheet1().c_pie_profit()
        sum_profit = int(round(df_profit['Unnamed: 13']).sum())
        profit_list = df_profit.values.tolist()
        rate_list = df_rate.values.tolist()
        # 修改文本
        body_shapes[2].text = pptText.Text.p2_text\
            .format(amount1=sum_profit, rate1=rate_list[1][3], rate2=rate_list[1][4], amount2=profit_list[0][0], rate3=rate_list[1][1],
                    amount3=profit_list[1][0], rate4=rate_list[1][1], amount4=profit_list[2][0], rate5=rate_list[1][2])
        body_shapes[2].text_frame.fit_text(font_family=u'SimSun-ExtB', max_size=14)
        # 添加图表
        TableChartMaker.Sheet1().c_barline_profit()
        p2_bar_left, p2_bar_top, p2_bar_width, p2_bar_height = Inches(0.05), Inches(1.9), Inches(8), Inches(4)  # 预设位置及大小
        slide.shapes.add_picture(self.pic_path + '大POS收益趋势图.jpg', p2_bar_left, p2_bar_top, p2_bar_width, p2_bar_height)
        TableChartMaker.Sheet1().c_pie_profit()
        p2_pie_left, p2_pie_top, p2_pie_width, p2_pie_height = Inches(8), Inches(1.9), Inches(4.5), Inches(4.2)
        slide.shapes.add_picture(self.pic_path + '3月收益分布.jpg', p2_pie_left, p2_pie_top, p2_pie_width, p2_pie_height)

        prs.save(self.save_path)
        print('第2页制作完成')

    def p3_operate(self):
        prs = Presentation(self.save_path)
        slide = prs.slides[2]
        body_shapes = slide.shapes

        left_title, top_title, width_title, height_title = Inches(4), Inches(1.4), Inches(6), Inches(1)  # 预设位置及大小
        textbox = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)  # left，top为相对位置，width，height为文本框大小
        textbox.text = '大pos收益与增长率              单位：万元'  # 文本框中文字
        textbox.text_frame.fit_text(font_family=u'SimSun-ExtB', max_size=20)

        df_table, r, c = TableChartMaker.Sheet1().t_bpos_profit()  # 获取表格行数、列数与数据
        rows_nums, cols_nums, left, top, width, height = r, c, Inches(1.5), Inches(2), Inches(10), Inches(5)
        table = body_shapes.add_table(rows_nums, cols_nums, left, top, width, height).table  # 添加表格，并取表格类
        fill_table(rows_nums, cols_nums, table, df_table)  # 填充数据
        alter_table_fontsize(table, 12)  # 修改表格字体大小

        prs.save(self.save_path)
        print('第3页制作完成')

    def p5_operate(self):
        prs = Presentation(self.save_path)
        slide = prs.slides[4]
        body_shapes = slide.shapes

        left, top, width, height = Inches(3.5), Inches(2), Inches(6), Inches(1)  # 预设位置及大小
        textbox = slide.shapes.add_textbox(left, top, width, height)  # left，top为相对位置，width，height为文本框大小
        textbox.text = '2020年4月大POS一级代理商给公司创造的总收益排名'  # 文本框中文字
        textbox.text_frame.fit_text(font_family=u'SimSun-ExtB', max_size=18)

        df_table, r, c = TableChartMaker.Sheet1().t_profit_ranking()  # 获取表格行数、列数与数据
        rows_nums, cols_nums, left, top, width, height = r, c, Inches(1.5), Inches(2.5), Inches(7), Inches(4)
        table = body_shapes.add_table(rows_nums, cols_nums, left, top, width, height).table  # 添加表格，并取表格类
        table.columns[3].width = Inches(4.5)  # 调整第4纵列宽度
        fill_table(rows_nums, cols_nums, table, df_table)  # 填充数据
        alter_table_fontsize(table, 12)  # 修改表格字体大小

        prs.save(self.save_path)
        print('第5页制作完成')

    def p9_operate(self):
        prs = Presentation(self.save_path)
        slide = prs.slides[8]
        body_shapes = slide.shapes
        # 获取文本所需数据
        pay_nums_type, nums_proportion = TableChartMaker.Sheet3().c_ring_pnums()
        pay_amo_type, amount_proportion = TableChartMaker.Sheet3().c_ring_pamount()
        # 修改文本
        body_shapes[2].text = pptText.Text.p9_text\
            .format(type1=pay_amo_type[0], proportion1=amount_proportion[0][0], type2=pay_amo_type[1], type3=pay_amo_type[2],
                    type4=pay_nums_type[0], proportion2=nums_proportion[0][0], type5=pay_nums_type[1], type6=pay_nums_type[2])
        body_shapes[2].text_frame.fit_text(font_family=u'SimSun-ExtB', max_size=14)
        # 添加图表
        TableChartMaker.Sheet3().c_ring_pamount()
        p9_left1, p9_top1, p9_width1, p9_height1 = Inches(1.5), Inches(2.5), Inches(5.1), Inches(4.5)  # 预设位置及大小
        slide.shapes.add_picture(self.pic_path + '3月各交易类型交易金额.jpg', p9_left1, p9_top1, p9_width1, p9_height1)
        TableChartMaker.Sheet3().c_ring_pnums()
        p9_left2, p9_top2, p9_width2, p9_height2 = Inches(7), Inches(2.5), Inches(5.1), Inches(4.5)  # 预设位置及大小
        slide.shapes.add_picture(self.pic_path + '3月各交易类型交易笔数.jpg', p9_left2, p9_top2, p9_width2, p9_height2)

        prs.save(self.save_path)
        print('第9页制作完成')
